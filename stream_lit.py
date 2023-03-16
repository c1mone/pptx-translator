import openai
import streamlit as st
from pptx import Presentation
from io import BytesIO

# Set up the Streamlit app
st.set_page_config(page_title="PowerPoint Translator")

# Set up the OpenAI API key input field
api_key = st.text_input("Enter your OpenAI API key:")

# Load the PowerPoint file
uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
if uploaded_file is not None:
    pr = Presentation(uploaded_file)

    # Count the total number of shapes in the PowerPoint file
    total_shapes = sum(len(slide.shapes) for slide in pr.slides)

    # Set up the progress bar
    progress_bar = st.progress(0)

    # Iterate through all the slides and replace the text with traditional Chinese
    shape_count = 0
    for slide in pr.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                original_text = shape.text
                if shape.text == "":
                    continue
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[{"role":"user","content":f"please translate \"{original_text}\" into traditional chinese"}],
                    temperature=0.2,
                    max_tokens=2048,
                    api_key=api_key
                )
                chinese_text = response["choices"][0]["message"]["content"]
                shape.text = chinese_text
                shape_count += 1
                progress_bar.progress(shape_count / total_shapes)

    buffer = BytesIO()
    pr.save(buffer)
    buffer.seek(0)
    # Download the updated PowerPoint file
    st.download_button(
        label="Download translated file",
        data=buffer,
        file_name="translated.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
